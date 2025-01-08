from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json

class PPTXGenerator:
    def __init__(self):
        self.prs = Presentation()
        
        self.alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }
        
    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def apply_text_style(self, text_frame, content):
        for item in content:
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = item["text"]
            
            if "style" in item:
                style = item["style"]
                
                if "font" in style:
                    font = style["font"]
                    if "name" in font:
                        run.font.name = font["name"]
                    if "size" in font:
                        run.font.size = Pt(font["size"])
                    if "bold" in font:
                        run.font.bold = font["bold"]
                    if "italic" in font:
                        run.font.italic = font["italic"]
                    if "color" in font:
                        rgb = self.hex_to_rgb(font["color"])
                        run.font.color.rgb = RGBColor(*rgb)
                
                if "paragraph" in style:
                    para_style = style["paragraph"]
                    if "alignment" in para_style:
                        paragraph.alignment = self.alignment_map.get(
                            para_style["alignment"].lower(),
                            PP_ALIGN.LEFT
                        )
                    if "lineSpacing" in para_style:
                        paragraph.line_spacing = para_style["lineSpacing"]
                    if "spaceBefore" in para_style:
                        paragraph.space_before = Pt(para_style["spaceBefore"])
                    if "spaceAfter" in para_style:
                        paragraph.space_after = Pt(para_style["spaceAfter"])
    
    def add_shape(self, slide, shape_data):
        pos = shape_data["position"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos["x"]),
            Inches(pos["y"]),
            Inches(pos["width"]),
            Inches(pos["height"])
        )
        
        if "style" in shape_data:
            style = shape_data["style"]
            if "fill" in style:
                fill = style["fill"]
                if "color" in fill:
                    rgb = self.hex_to_rgb(fill["color"])
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*rgb)
            
            if "line" in style:
                line = style["line"]
                if "color" in line:
                    rgb = self.hex_to_rgb(line["color"])
                    shape.line.color.rgb = RGBColor(*rgb)
                if "width" in line:
                    shape.line.width = Pt(line["width"])
    
    def add_picture(self, slide, picture_data):
       
        pos = picture_data["position"]
        picture = slide.shapes.add_picture(
            picture_data["path"],
            Inches(pos["x"]),
            Inches(pos["y"]),
            Inches(pos["width"]),
            Inches(pos["height"])
        )
        
        if "crop" in picture_data:
            crop = picture_data["crop"]
            picture.crop_left = crop.get("left", 0)
            picture.crop_top = crop.get("top", 0)
            picture.crop_right = crop.get("right", 0)
            picture.crop_bottom = crop.get("bottom", 0)
    
    def add_table(self, slide, table_data):
      
        pos = table_data["position"]
        rows = table_data["rows"]
        cols = table_data["cols"]
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(pos["x"]),
            Inches(pos["y"]),
            Inches(pos["width"]),
            Inches(pos["height"])
        ).table
        
        for row_idx, row in enumerate(table_data["data"]):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)
        
        if "style" in table_data:
            style = table_data["style"]
            if "cells" in style:
                for cell_style in style["cells"]:
                    row = cell_style["row"]
                    col = cell_style["col"]
                    cell = table.cell(row, col)
                    
                    if "fill" in cell_style:
                        fill = cell_style["fill"]
                        if "color" in fill:
                            rgb = self.hex_to_rgb(fill["color"])
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(*rgb)
                    
                    if "text" in cell_style:
                        text_frame = cell.text_frame
                        text_frame.clear()
                        self.apply_text_style(text_frame, [cell_style["text"]])

    def create_slide(self, slide_data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        for element in slide_data["elements"]:
            element_type = element["type"]
            
            if element_type == "text":
                textbox = slide.shapes.add_textbox(
                    Inches(element["position"]["x"]),
                    Inches(element["position"]["y"]),
                    Inches(element["position"]["width"]),
                    Inches(element["position"]["height"])
                )
                self.apply_text_style(textbox.text_frame, element["content"])
            
            elif element_type == "shape":
                self.add_shape(slide, element)
            
            elif element_type == "picture":
                self.add_picture(slide, element)
            
            elif element_type == "table":
                self.add_table(slide, element)
    
    def generate_presentation(self, json_data):
        if "metadata" in json_data["presentation"]:
            core_properties = self.prs.core_properties
            metadata = json_data["presentation"]["metadata"]
            
            if "title" in metadata:
                core_properties.title = metadata["title"]
            if "author" in metadata:
                core_properties.author = metadata["author"]
            if "subject" in metadata:
                core_properties.subject = metadata["subject"]
        
        for slide_data in json_data["presentation"]["slides"]:
            self.create_slide(slide_data)
    
    def save(self, filename):
        self.prs.save(filename)

def create_presentation_from_json(json_file_path, output_pptx_path):
    with open(json_file_path, 'r') as f:
        json_data = json.load(f)
    
    generator = PPTXGenerator()
    generator.generate_presentation(json_data)
    generator.save(output_pptx_path)

if __name__ == "__main__":
    create_presentation_from_json("ml_presentation.json", "output.pptx")